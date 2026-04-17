from __future__ import annotations

import hashlib
import json
import re
import uuid
from pathlib import Path

from app.config import CHUNK_OVERLAP_CHARS, ENABLE_DOCLING, MAX_CHUNK_CHARS, SILVER_NOTES_DIR
from app.schemas import DocumentUnit, FinalSolutionRecord, TeacherNoteRecord

SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx"}
PROBLEM_HEADER_PATTERN = re.compile(r"(?m)^(?P<problem_id>\d+\.\d+)\s+(?P<title>.+)$")


def classify_source_family(path: Path) -> str:
    name = path.name.lower()
    if "final" in name:
        return "final_exam"
    if "solution" in name:
        return "solution_gold"
    if "problem" in name or "exercise" in name or "question" in name:
        return "problem_gold"
    return "theory_gold"


def extract_chapter(path: Path) -> str | None:
    match = re.search(r"chapter\s*(\d+)", path.stem, re.IGNORECASE)
    if match:
        return f"chapter{match.group(1)}"
    return None


def extract_problem_id(text: str) -> str | None:
    match = re.search(r"\b(\d+\.\d+)\b", text)
    if match:
        return match.group(1)
    return None


def pair_key(chapter: str | None, problem_id: str | None) -> str | None:
    if not chapter or not problem_id:
        return None
    return f"{chapter}_problem_{problem_id}"


def _normalize_text(text: str) -> str:
    lines = [line.strip() for line in text.replace("\x00", " ").splitlines()]
    return "\n".join(line for line in lines if line)


def split_text_segments(text: str) -> list[str]:
    text = _normalize_text(text)
    if not text:
        return []
    if len(text) <= MAX_CHUNK_CHARS:
        return [text]

    chunks: list[str] = []
    start = 0
    text_length = len(text)
    while start < text_length:
        end = min(text_length, start + MAX_CHUNK_CHARS)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= text_length:
            break
        start = max(end - CHUNK_OVERLAP_CHARS, start + 1)
    return chunks


def extract_pdf_pages(path: Path) -> list[tuple[int, str]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Missing dependency: pypdf") from exc

    reader = PdfReader(str(path))
    pages: list[tuple[int, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        normalized = _normalize_text(text)
        if normalized:
            pages.append((idx, normalized))
    return pages


def extract_ppt_slides(path: Path) -> list[tuple[int, str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Missing dependency: python-pptx") from exc

    presentation = Presentation(str(path))
    slides: list[tuple[int, str]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                texts.append(text)
        normalized = _normalize_text("\n".join(texts))
        if normalized:
            slides.append((idx, normalized))
    return slides


def extract_docling_units(path: Path) -> list[tuple[int, str]]:
    if not ENABLE_DOCLING:
        return []
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        return []
    try:
        converter = DocumentConverter()
        result = converter.convert(str(path))
        markdown_text = _normalize_text(result.document.export_to_markdown())
    except Exception:
        return []
    if not markdown_text:
        return []
    blocks = [block.strip() for block in markdown_text.split("\n## ") if block.strip()]
    if not blocks:
        return [(1, markdown_text)]
    units: list[tuple[int, str]] = []
    for index, block in enumerate(blocks, start=1):
        units.append((index, block if index == 1 else f"## {block}"))
    return units


def extract_units(path: Path) -> list[tuple[int, str]]:
    docling_units = extract_docling_units(path)
    if docling_units:
        return docling_units

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_pages(path)
    if suffix in {".ppt", ".pptx"}:
        return extract_ppt_slides(path)
    raise RuntimeError(f"Unsupported file type: {suffix}")


def build_document_units(path: Path, *, include_finals: bool = False) -> list[DocumentUnit]:
    source_family = classify_source_family(path)
    if source_family == "final_exam" and not include_finals:
        return []

    if source_family in {"problem_gold", "solution_gold"}:
        units = build_problem_solution_units(path, source_family)
        if units:
            return units
    return build_text_units(path, source_family)


def build_text_units(path: Path, source_family: str | None = None) -> list[DocumentUnit]:
    source_family = source_family or classify_source_family(path)
    chapter = extract_chapter(path)
    source_type = path.suffix.lower().lstrip(".")
    unit_kind = "concept_note" if source_family == "theory_gold" else "text"
    is_answer_like = source_family == "solution_gold"
    records: list[DocumentUnit] = []

    for page_or_slide, unit_text in extract_units(path):
        title = first_heading(unit_text)
        for chunk_index, chunk_text in enumerate(split_text_segments(unit_text)):
            records.append(
                DocumentUnit(
                    point_id=stable_point_id(path, source_family, page_or_slide, chunk_index, title or ""),
                    file_path=str(path.resolve()),
                    source_type=source_type,
                    source_family=source_family,
                    source_quality="gold",
                    chapter=chapter,
                    title=title,
                    unit_kind=unit_kind,
                    page_or_slide=page_or_slide,
                    chunk_index=chunk_index,
                    is_answer_like=is_answer_like,
                    text=chunk_text,
                    excerpt=chunk_text[:240],
                    metadata={
                        "page_or_slide": page_or_slide,
                        "chapter": chapter,
                        "source_family": source_family,
                        "title": title,
                    },
                )
            )
    return records


def build_problem_solution_units(path: Path, source_family: str) -> list[DocumentUnit]:
    pages = extract_units(path)
    if not pages:
        return []
    chapter = extract_chapter(path)
    source_type = path.suffix.lower().lstrip(".")
    full_text, page_markers = join_with_page_markers(pages)
    matches = list(PROBLEM_HEADER_PATTERN.finditer(full_text))
    if not matches:
        return []

    units: list[DocumentUnit] = []
    for index, match in enumerate(matches):
        start = match.start()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(full_text)
        raw_segment = _normalize_text(re.sub(r"\[\[PAGE_\d+\]\]", " ", full_text[start:end]))
        problem_id = match.group("problem_id")
        title = _normalize_text(match.group("title"))
        page_or_slide = page_for_offset(start, page_markers)
        key = pair_key(chapter, problem_id)
        segment_kind = "solution" if source_family == "solution_gold" else "problem"

        for chunk_index, chunk_text in enumerate(split_text_segments(raw_segment)):
            units.append(
                DocumentUnit(
                    point_id=stable_point_id(path, source_family, problem_id, chunk_index, title),
                    file_path=str(path.resolve()),
                    source_type=source_type,
                    source_family=source_family,
                    source_quality="gold",
                    chapter=chapter,
                    problem_id=problem_id,
                    pair_key=key,
                    title=title,
                    unit_kind=segment_kind,
                    page_or_slide=page_or_slide,
                    chunk_index=chunk_index,
                    is_answer_like=source_family == "solution_gold",
                    text=chunk_text,
                    excerpt=chunk_text[:240],
                    metadata={
                        "page_or_slide": page_or_slide,
                        "chapter": chapter,
                        "problem_id": problem_id,
                        "pair_key": key,
                        "title": title,
                        "source_family": source_family,
                    },
                )
            )
    return units


def load_teacher_notes(include_silver: bool = True) -> list[DocumentUnit]:
    if not include_silver:
        return []

    notes: list[DocumentUnit] = []
    for path in sorted(SILVER_NOTES_DIR.glob("*.json")):
        try:
            payload = json.loads(path.read_text(encoding="utf-8-sig"))
        except Exception:
            continue
        records = payload if isinstance(payload, list) else [payload]
        for index, item in enumerate(records):
            if is_final_solution_payload(item):
                notes.extend(build_final_solution_units(path, item, index))
                continue
            try:
                note = TeacherNoteRecord.model_validate(item)
            except Exception:
                continue
            notes.append(
                DocumentUnit(
                    point_id=stable_point_id(path, "notes_silver", note.note_id, index, note.title),
                    file_path=str(path.resolve()),
                    source_type="json",
                    source_family="notes_silver",
                    source_quality="silver",
                    chapter=note.chapter,
                    pair_key=note.pair_key,
                    title=note.title,
                    unit_kind=note.note_type,
                    is_answer_like=False,
                    text=_normalize_text(note.text),
                    excerpt=_normalize_text(note.text)[:240],
                    metadata={
                        "pair_key": note.pair_key,
                        "chapter": note.chapter,
                        "verification_status": note.verification_status,
                        "teacher_model": note.teacher_model,
                        **note.metadata,
                    },
                )
            )
    return notes


def is_final_solution_payload(item: object) -> bool:
    if not isinstance(item, dict):
        return False
    required_keys = {"exam_id", "question_id", "full_solution", "final_answer"}
    return required_keys.issubset(item)


def build_final_solution_units(path: Path, item: dict[str, object], index: int) -> list[DocumentUnit]:
    try:
        record = FinalSolutionRecord.model_validate(item)
    except Exception:
        return []

    pair = final_pair_key(record.exam_id, record.question_id)
    combined_text = _normalize_text(
        "\n".join(
            [
                f"Exam: {record.exam_id}",
                f"Question: {record.question_id}",
                f"Title: {record.title}",
                f"Question Text:\n{record.question_text}",
                f"Full Solution:\n{record.full_solution}",
                f"Final Answer:\n{record.final_answer}",
                f"Key Formulas: {', '.join(record.key_formulas)}" if record.key_formulas else "",
                f"Method Tags: {', '.join(record.method_tags)}" if record.method_tags else "",
            ]
        )
    )
    units: list[DocumentUnit] = []
    for chunk_index, chunk_text in enumerate(split_text_segments(combined_text)):
        units.append(
            DocumentUnit(
                point_id=stable_point_id(path, "final_solution_silver", record.solution_id, index, chunk_index),
                file_path=str(path.resolve()),
                source_type="json",
                source_family="final_solution_silver",
                source_quality="silver",
                problem_id=record.question_id,
                pair_key=pair,
                title=record.title,
                unit_kind="final_solution",
                chunk_index=chunk_index,
                is_answer_like=True,
                text=chunk_text,
                excerpt=chunk_text[:240],
                metadata={
                    "exam_id": record.exam_id,
                    "question_id": record.question_id,
                    "pair_key": pair,
                    "verification_status": record.verification_status,
                    "teacher_model": record.teacher_model,
                    "key_formulas": record.key_formulas,
                    "method_tags": record.method_tags,
                    **record.metadata,
                },
            )
        )
    return units


def final_pair_key(exam_id: str, question_id: str) -> str:
    normalized_exam = re.sub(r"[^a-zA-Z0-9]+", "_", exam_id).strip("_")
    normalized_question = re.sub(r"[^a-zA-Z0-9]+", "_", question_id).strip("_")
    return f"{normalized_exam}_{normalized_question}"


def join_with_page_markers(pages: list[tuple[int, str]]) -> tuple[str, list[tuple[int, int]]]:
    parts: list[str] = []
    markers: list[tuple[int, int]] = []
    offset = 0
    for page_number, text in pages:
        marker = f"\n[[PAGE_{page_number}]]\n"
        parts.append(marker)
        offset += len(marker)
        markers.append((offset, page_number))
        parts.append(text)
        offset += len(text)
    return "".join(parts), markers


def page_for_offset(offset: int, page_markers: list[tuple[int, int]]) -> int | None:
    page_number: int | None = None
    for marker_offset, candidate in page_markers:
        if marker_offset <= offset:
            page_number = candidate
        else:
            break
    return page_number or (page_markers[0][1] if page_markers else None)


def first_heading(text: str) -> str | None:
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        if len(line) <= 140:
            return line
        return line[:137] + "..."
    return None


def stable_point_id(path: Path, *parts: object) -> str:
    raw = "::".join([str(path.resolve()), *[str(part) for part in parts]])
    return str(uuid.uuid5(uuid.NAMESPACE_URL, raw))
